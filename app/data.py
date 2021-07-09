import cv2
import numpy as np
import pytesseract
from pytesseract import Output
import nltk
from nltk import word_tokenize
nltk.download('punkt')
from nltk import ngrams, FreqDist
import pandas as pd
import matplotlib.pyplot as plt
nltk.download('stopwords') 
from nltk.corpus import stopwords
nltk.download('wordnet')
from nltk.stem import WordNetLemmatizer
from sklearn.feature_extraction.text import TfidfVectorizer
from pdf2image import convert_from_path
# from IPython.display import Image
from pptx import Presentation
import glob
import docx
import os
from sklearn.metrics.pairwise import cosine_similarity
import cx_Oracle


def get_grayscale(image):
    return cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

def getImageText(filename):
    img = cv2.imread(filename)
    grayscale = get_grayscale(img)
    custom_config = r'-l eng --psm 6'
    text = pytesseract.image_to_string(grayscale, config=custom_config)
    return text

def getPDFText(filename):
    pages = convert_from_path(filename, 500)
    pagenumber = 1
    page_paths = []
    for page in pages:
        filename = "page"+ str(pagenumber) + ".jpg"
        pagenumber += 1
        page_paths.append(filename)
        page.save(filename, 'JPEG')
    images = []
    for page in page_paths:
        images.append(cv2.imread(page))
    for image in images:
        image = get_grayscale(image)
    custom_config = r'-l eng --psm 6'
    text = ""
    for grayscale in images:
        text += pytesseract.image_to_string(grayscale, config=custom_config)
    return text

def getPPTXText(filename):
    for eachfile in glob.glob(filename):
        prs = Presentation(eachfile)
        fullText = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    fullText.append(shape.text)
    return '\n'.join(fullText)

def getDOCXText(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def detectFileType(filename):
    extension = filename[filename.rfind(".")+1:]
    return extension

def getFileName(filepath):
    identifier = "/"
    filename = filepath[filepath.rfind(identifier)+1:]
    return filename

def preprocessing(text):
    words = word_tokenize(text)
    updated_text = [word.lower() for word in words]
    final_text = [x for x in updated_text if x.isalpha()]
    stop_words = stopwords.words('english')
    final_text = [word for word in final_text if word not in stop_words]
    lemmatizer = WordNetLemmatizer()
    for word in final_text:
        word = lemmatizer.lemmatize(word)
    final_string = " ".join(final_text)
    return final_string

def createFramework(connection):
    try:
        con = cx_Oracle.connect(connection)
        cursor = con.cursor()
        cursor.execute("CREATE TABLE FILENAMES (fileID int NOT NULL PRIMARY KEY, filepath varchar(255), filename varchar(255))")

        print("File Name Table created successfully.")
        
        cursor.execute("CREATE TABLE TEXT (fileID int, text LONG, FOREIGN KEY (fileID) REFERENCES FILENAMES(fileID))")
        
        print("Text Table created successfully.")
        
        cursor.execute("CREATE TABLE FILES (fileID int, filedata BLOB, FOREIGN KEY (fileID) REFERENCES FILENAMES(fileID))")
        
        print("Files Table created successfully.")
        
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close() 

def deleteFramework(connection):
    try:
        con = cx_Oracle.connect(connection)
        cursor = con.cursor()
        
        cursor.execute("DROP TABLE TEXT")
        
        print("Text Table dropped successfully.")
        
        cursor.execute("DROP TABLE FILES")
        
        print("Files Table dropped successfully")
        
        cursor.execute("DROP TABLE FILENAMES")

        print("File Name Table dropped successfully.")

    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close() 

def convertToBinaryData(filename):
    with open(filename, 'rb') as file:
        binaryData = file.read()
    return binaryData

def write_file(data, filename):
    with open(filename, 'wb') as file:
        file.write(data)

def stringToBinary(string):
    final_string = ''.join(format(i, '08b') for i in bytearray(string, encoding ='utf-8'))
    return final_string

def binaryToString(binary):
    return ''.join(chr(int(binary[i*8:i*8+8],2)) for i in range(len(binary)//8))

def getText(filepath):
    file_type = detectFileType(filepath)
    string_text = ""
    if file_type == "jpg":
        string_text = getImageText(filepath)
    elif file_type == "pdf":
        string_text = getPDFText(filepath)
    elif file_type == "pptx":
        string_text = getPPTXText(filepath)
    elif file_type == "docx":
        string_text = getDOCXText(filepath)
    return preprocessing(string_text)
    
def loadFileName(filepath):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        filename = getFileName(filepath)
        
        
        cursor.execute("select max(fileID) from FILENAMES") 
        
        for maxID in cursor:
            maxID = maxID[0]
            if maxID == None:
                fileID = 1
            else:
                fileID = maxID + 1

        sql = ('insert into FILENAMES(fileID, filepath, filename) values(:fileID,:filepath,:filename)')
        
        cursor.execute(sql, [fileID, filepath, filename])
        
        con.commit()
        
        print("Successfuly loaded file name into FILENAMES")
    
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def getFileID(file_path):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        sql = ('select * from FILENAMES')
        
        cursor.execute(sql)
        
        for fileID, filepath, filename in cursor:
            if filepath == file_path:
                return fileID
    
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def loadFileData(filepath):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        fileID = getFileID(filepath)
        
        filedata = convertToBinaryData(filepath)
        
        sql = ('insert into FILES(fileID, filedata) values(:fileID,:filedata)')
        
        cursor.execute(sql, [fileID, filedata])
        
        con.commit()
        
        print("Successfuly loaded file data into FILES")
    
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def loadInText(filepath):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        fileID = getFileID(filepath)
    
        fileText = getText(filepath)
        
        text = stringToBinary(fileText)
        
        sql = ('insert into TEXT(fileID, text) values(:fileID,:text)')
        
        cursor.execute(sql, [fileID, text])
        
        con.commit()
        
        print("Successfuly loaded file text into TEXT")
    
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def loadInFile(filepath):
    loadFileName(filepath)
    loadFileData(filepath)
    loadInText(filepath)
    print("Successfully loaded in file into database")

def downloadFile(file_ID):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        file_name = ""
        
        sql = ('select * from FILENAMES')
        
        cursor.execute(sql)
        
        for fileID, filepath, filename in cursor:
            if fileID == file_ID:
                file_name = "\\" + filename
                
        sql = ('select * from FILES')
        
        cursor.execute(sql)
        
        for fileID, filedata in cursor:
            if fileID == file_ID:
                filepath = r"C:\Users\pradi\Documents\NPS2021\files" + file_name
                print("Opening " + str(filepath) + "...")
                file_data = filedata.read()
                write_file(file_data, filepath)
                os.system("start "+filepath)
                print("Opened file!")
    
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def filenameFromID(file_ID):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        file_name = ""
        
        sql = ('select * from FILENAMES')
        
        cursor.execute(sql)
        
        for fileID, filepath, filename in cursor:
            if fileID == file_ID:
                file_name = filename
                return file_name
            
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def vectorizeData(text):
    vectorizer = TfidfVectorizer()
    vectors = vectorizer.fit_transform([text[file] for file in text])
    dense = vectors.todense()
    
    cosine_sim = cosine_similarity(dense, dense)
    
    return cosine_sim

def max_score(score_list):
    max_score = 0
    index = 0
    max_index = 0
    for score in score_list:
        if score >= max_score:
            max_score = score
            max_index = index
        index += 1
  
    return max_index, max_score

def matchQuery(query):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        sql = ('select * from TEXT')
        
        cursor.execute(sql)
        
        text = {}
        
        for fileID, fileText in cursor:
            text[fileID] = binaryToString(fileText)
                
        text["query"] = preprocessing(query)
        
        similarityMatrix = vectorizeData(text)
        
        matrix_length = len(similarityMatrix)
        
        query_scores = similarityMatrix[matrix_length-1]
        query_scores = query_scores[0:matrix_length-1]
        
        topScores = {}
        
        max_index1, topScores[max_index1] = max_score(query_scores)
        query_scores[max_index1] = -1
        max_index2, topScores[max_index2] = max_score(query_scores)
        query_scores[max_index2] = -1
        max_index3, topScores[max_index3] = max_score(query_scores)
        query_scores[max_index3] = -1
        
        return topScores
        
        
    
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def openDocument(query):
    try: 
        con = cx_Oracle.connect('YOSHI/nps2021@localhost')
        cursor = con.cursor()
        
        sql = ('select * from FILENAMES')
        
        cursor.execute(sql)
        
        fileID = 0
        
        for fileID, filepath, filename in cursor:
            if filename == query:
                downloadFile(fileID)
        
    except cx_Oracle.DatabaseError as e:
        print("There was a problem with Oracle", e)


    finally: 
        if cursor:
            cursor.close()
        if con:
            con.close()

def getQuery(query):
    
    topScores = matchQuery(query)

    fileNames = {}
    
    for fileID in topScores:
        filename = filenameFromID(fileID+1)
        fileNames[filename] = topScores[fileID]
    
    return filesNames








